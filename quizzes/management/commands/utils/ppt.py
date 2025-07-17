from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE
from pptx.shapes.placeholder import Shape
from pptx.slide import Slide
from pptx.util import Inches

from quizzes.management.commands.utils.html import get_sorted_themes


def create_ppt(data, output_file):
    left = top = width = height = Inches(0.75)
    prs = Presentation()

    # Generate index slide
    slide_layout = prs.slide_layouts[1]
    home_slide = prs.slides.add_slide(slide_layout)
    title = home_slide.shapes.title
    title.text = output_file

    themes = get_sorted_themes(data)
    theme_to_first_question_slide = {}

    # Generate slides for each question
    question_slides = []
    for row in data:
        slide = prs.slides.add_slide(slide_layout)
        # title
        title = slide.shapes.title
        title.text = row["theme"]
        # question
        body_shape: Shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        tf.text = row["question"]

        if row["theme"] not in theme_to_first_question_slide:
            theme_to_first_question_slide[row["theme"]] = slide
        question_slides.append(slide)

        slide = prs.slides.add_slide(slide_layout)
        # title
        title = slide.shapes.title
        title.text = row["theme"]
        # answer
        body_shape: Shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = row["answer"]
        # footer with xP
        footer = slide.shapes.add_textbox(left * 2, prs.slide_height - top, prs.slide_width - 4 * width, height)
        footer.text_frame.text = f"xP: {row['xP']}"
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("c0c0c0")

    # Add hyperlinks from the index slide to the first question slide of each theme
    for i, theme in enumerate(themes, 1):
        theme_slide: Slide = theme_to_first_question_slide[theme]
        index_shape: Shape = home_slide.shapes.placeholders[1]
        rId = home_slide.part.relate_to(theme_slide.part, RELATIONSHIP_TYPE.SLIDE)
        p = index_shape.text_frame.add_paragraph()
        r = p.add_run()
        r.text = theme
        rPr = r._r.get_or_add_rPr()
        hlinkClick = rPr.add_hlinkClick(rId)
        hlinkClick.set("action", "ppaction://hlinksldjump")

    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        # all slides have button back to the index slide
        home_button = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ACTION_BUTTON_HOME, prs.slide_width - left, prs.slide_height - top, width, height
        )
        home_button.click_action.target_slide = home_slide

    return prs
